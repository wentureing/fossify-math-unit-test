#!/usr/bin/env python3
"""根据 ppt设计方案-codebuddy.md 生成 PPTX 文件"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap
from pptx.oxml import parse_xml

# ======================== 颜色常量 ========================
DARK_BLUE = RGBColor(0x1A, 0x23, 0x7E)
MID_BLUE = RGBColor(0x28, 0x35, 0x93)
ACCENT_BLUE = RGBColor(0x19, 0x76, 0xD2)
LIGHT_BLUE = RGBColor(0xE3, 0xF2, 0xFD)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
DARK_GRAY = RGBColor(0x42, 0x42, 0x42)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
GREEN = RGBColor(0x43, 0xA0, 0x47)
ORANGE = RGBColor(0xFB, 0x8C, 0x00)
RED = RGBColor(0xC6, 0x28, 0x28)
PURPLE = RGBColor(0x7B, 0x1F, 0xA2)
TEAL = RGBColor(0x00, 0x89, 0x6B)

# ======================== 工具函数 ========================
def set_shape_fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color

def set_shape_no_fill(shape):
    shape.fill.background()

def add_text_box(slide, left, top, width, height, text, font_size=14,
                 bold=False, color=BLACK, align=PP_ALIGN.LEFT,
                 font_name="Microsoft YaHei", fill_color=None):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    if fill_color:
        set_shape_fill(txBox, fill_color)
    return txBox

def add_title_shape(slide, text, left, top, width, height,
                    font_size=28, color=WHITE, bold=True,
                    fill_color=DARK_BLUE, align=PP_ALIGN.CENTER):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    set_shape_fill(shape, fill_color)
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = "Microsoft YaHei"
    p.alignment = align
    tf.margin_left = Pt(10)
    tf.margin_right = Pt(10)
    tf.margin_top = Pt(6)
    tf.margin_bottom = Pt(6)
    return shape

def add_table(slide, rows, cols, left, top, width, height, data,
              col_widths=None, header_color=DARK_BLUE, header_font_size=11,
              cell_font_size=10, row_heights=None):
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w
    
    for i in range(rows):
        if row_heights and i < len(row_heights):
            table.rows[i].height = row_heights[i]
        for j in range(cols):
            cell = table.cell(i, j)
            cell.text = str(data[i][j]) if i < len(data) and j < len(data[i]) else ""
            para = cell.text_frame.paragraphs[0]
            para.font.size = Pt(header_font_size if i == 0 else cell_font_size)
            para.font.name = "Microsoft YaHei"
            para.font.bold = (i == 0)
            para.alignment = PP_ALIGN.CENTER
            
            if i == 0:
                set_shape_fill(cell, header_color)
                para.font.color.rgb = WHITE
            else:
                # 斑马纹
                if i % 2 == 0:
                    set_shape_fill(cell, LIGHT_GRAY)
                else:
                    set_shape_fill(cell, WHITE)
                para.font.color.rgb = BLACK
    
    return table

def add_bullet_points(slide, left, top, width, height, items,
                      font_size=12, color=BLACK, bullet_char="•"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for idx, item in enumerate(items):
        if idx == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"{bullet_char} {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Microsoft YaHei"
        p.space_after = Pt(4)
    return txBox

# ======================== PPT 页面生成 ========================

def create_cover(prs):
    """P01 封面"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局
    
    # 背景
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    set_shape_fill(bg, DARK_BLUE)
    
    # 装饰条
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(3.4), Inches(8), Pt(3))
    set_shape_fill(bar, WHITE)
    
    # 主标题
    add_text_box(slide, Inches(0.5), Inches(1.5), Inches(9), Inches(1),
                 "Fossify Calculator", 44, True, WHITE, PP_ALIGN.CENTER)
    
    # 副标题
    add_text_box(slide, Inches(0.5), Inches(2.3), Inches(9), Inches(0.6),
                 "Android 开源计算器应用测试报告", 26, False, RGBColor(0xBB, 0xDE, 0xFB), PP_ALIGN.CENTER)
    
    # 测试维度
    add_text_box(slide, Inches(0.5), Inches(3.7), Inches(9), Inches(0.6),
                 "单元测试 / 集成测试 / 性能测试 / 覆盖率分析", 18, False, WHITE, PP_ALIGN.CENTER)
    
    # 课程信息
    info_items = [
        ("课程名称", "软件测试"),
        ("汇报人", "XXX"),
        ("学号", "XXXXXX"),
        ("指导教师", "XXX"),
        ("日期", "2026 年 X 月 X 日"),
    ]
    y_start = 4.6
    for label, value in info_items:
        add_text_box(slide, Inches(3), Inches(y_start), Inches(2), Inches(0.35),
                     f"{label}：{value}", 13, False, WHITE, PP_ALIGN.LEFT)
        y_start += 0.32


def create_background(prs):
    """P02 项目背景"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 标题栏
    add_title_shape(slide, "一、项目背景与测试目标", 0, 0, prs.slide_width, Inches(0.65))
    
    # 左侧：应用简介表格
    app_data = [
        ["维度", "描述"],
        ["应用名称", "Fossify Calculator"],
        ["应用类型", "Android 开源计算器"],
        ["开源协议", "完全开源 (GPL v3)"],
        ["离线支持", "无需网络权限"],
        ["隐私保护", "不收集用户数据"],
    ]
    add_table(slide, len(app_data), 2, Inches(0.3), Inches(0.85), Inches(4.2), Inches(2.3), app_data,
              col_widths=[Inches(1.4), Inches(2.8)], header_font_size=11, cell_font_size=10)
    
    # 右侧：核心功能
    func_data = [
        ["模块", "功能"],
        ["🧮 基础运算", "加减乘除、小数、负数运算"],
        ["🔬 高级运算", "幂运算、开方、百分比"],
        ["📏 单位转换", "9大类109种单位相互转换"],
        ["📊 计算历史", "历史记录保存、查看、清空"],
        ["⚙️ 自定义设置", "振动开关、屏幕常亮、主题颜色"],
        ["🖼️ 桌面小部件", "App Widget 快捷计算"],
    ]
    add_table(slide, len(func_data), 2, Inches(4.7), Inches(0.85), Inches(4.6), Inches(2.3), func_data,
              col_widths=[Inches(1.6), Inches(3.0)], header_font_size=11, cell_font_size=10)
    
    # 测试目标表格（左下）
    goal_data = [
        ["测试维度", "目标", "意义"],
        ["功能正确性", "验证核心计算逻辑准确性", "保证用户体验基础"],
        ["代码质量", "提升代码可维护性", "降低后期维护成本"],
        ["性能稳定性", "确保流畅运行体验", "提升用户满意度"],
        ["安全性验证", "验证隐私合规性", "保护用户数据安全"],
    ]
    add_table(slide, len(goal_data), 3, Inches(0.3), Inches(3.3), Inches(9), Inches(1.75), goal_data,
              col_widths=[Inches(1.6), Inches(3.5), Inches(3.9)], header_font_size=11, cell_font_size=9)
    
    # 时间线说明（右下）
    timeline_items = [
        "准备阶段: 源码获取与环境配置 (3d)",
        "       项目编译与初步运行 (3d)",
        "测试实施: 单元测试编写与执行 (7d)",
        "       集成测试编写与执行 (5d)",
        "分析总结: 覆盖率分析与报告生成 (3d)",
        "       性能测试与瓶颈分析 (3d)",
    ]
    add_bullet_points(slide, Inches(0.5), Inches(5.15), Inches(9), Inches(1.5),
                      timeline_items, 10, DARK_GRAY)


def create_architecture(prs):
    """P03 系统架构"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    add_title_shape(slide, "二、系统架构与代码规模分析", 0, 0, prs.slide_width, Inches(0.65))
    
    # 技术栈标签云（左侧上）
    tech_data = [
        ["技术", "用途", "版本要求"],
        ["Kotlin", "主要开发语言", "1.9+"],
        ["Android SDK", "目标平台", "API 33"],
        ["Jetpack Compose", "设置页面 UI", "BOM 2024.x"],
        ["Room Database", "本地持久化存储", "2.x"],
        ["EvalEx", "数学表达式解析库", "最新版"],
        ["BigDecimal", "高精度数值计算", "DECIMAL128"],
        ["JUnit 4 / MockK", "单元测试框架", "4.13.2 / 1.13.x"],
        ["Robolectric / Espresso", "组件/UI 测试", "4.x / 3.5.x"],
        ["Jacoco", "代码覆盖率工具", "0.8.x"],
    ]
    add_table(slide, len(tech_data), 3, Inches(0.25), Inches(0.78), Inches(4.5), Inches(2.95), tech_data,
              col_widths=[Inches(1.5), Inches(2.0), Inches(1.0)], header_font_size=10, cell_font_size=9)
    
    # 代码规模统计表（右侧上）
    scale_data = [
        ["指标", "数量", "说明"],
        ["Package 数量", "10 个", "org.fossify.math.* 下 10 个子包"],
        ["Kotlin 源文件", "35 个", "app/src/main/kotlin 下的 .kt 文件"],
        ["Class 数量", "32+ 个", "含 Activity、类、接口、单例、数据类"],
        ["Method 数量", "~240+ 个", "含重写方法、私有方法、扩展函数"],
        ["代码总行数", "~3500+ 行", "不含自动生成的 BuildConfig 等"],
        ["Activity 数量", "7 个", "主界面、设置、单位转换等"],
    ]
    add_table(slide, len(scale_data), 3, Inches(4.85), Inches(0.78), Inches(4.55), Inches(2.95), scale_data,
              col_widths=[Inches(1.5), Inches(1.0), Inches(2.05)], header_font_size=10, cell_font_size=9)
    
    # 包结构树形图（下方）
    pkg_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     Inches(0.25), Inches(3.85), Inches(9.2), Inches(2.2))
    set_shape_fill(pkg_box, RGBColor(0xE8, 0xEA, 0xF6))
    
    add_text_box(slide, Inches(0.4), Inches(3.92), Inches(3), Inches(0.3),
                 "包结构 org.fossify.math.*", 12, True, DARK_BLUE, PP_ALIGN.LEFT)
    
    pkg_items = [
        "activities     → 7个Activity (MainActivity/Settings/UnitConverter...)",
        "adapters       → 2个RecyclerView适配器",
        "compose        → 1个Compose组件 (SettingsScreen)",
        "databases      → Room数据库单例 + DAO接口",
        "dialogs        → 历史对话框",
        "extensions     → 6个扩展函数",
        "helpers        → 核心业务逻辑 ★",
        "   ├── CalculatorImpl    → 核心计算引擎 (~20方法)",
        "   ├── NumberFormatHelper→ 数字格式化工具",
        "   └── converters       → 9种单位转换器 (Length/Area/Volume/Mass/Temp/Time/Speed/Pressure/Energy)",
        "interfaces      → DAO接口定义",
        "models         → 数据模型 (History Entity / ConverterUnitsState)",
        "receivers      → 语言切换广播接收器",
        "views          → 自定义转换视图 (ConverterView)",
    ]
    y = 4.22
    for item in pkg_items:
        add_text_box(slide, Inches(0.45), Inches(y), Inches(8.8), Inches(0.18),
                     item, 9, False, DARK_GRAY, PP_ALIGN.LEFT)
        y += 0.17


def create_class_diagram(prs):
    """P04 类关系图"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    add_title_shape(slide, "三、核心类关系图", 0, 0, prs.slide_width, Inches(0.65))
    
    # 左侧：核心依赖关系图（用文本框模拟）
    diagram_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                         Inches(0.2), Inches(0.78), Inches(5.6), Inches(4.9))
    set_shape_fill(diagram_box, RGBColor(0xFA, 0xFA, 0xFA))
    
    add_text_box(slide, Inches(0.35), Inches(0.86), Inches(5), Inches(0.28),
                 "核心类关系简化视图", 12, True, ACCENT_BLUE, PP_ALIGN.LEFT)
    
    # 类关系文字描述
    class_relations = [
        ("┌─────────────────────┐", ""),
        ("│   MainActivity      │ ◄── 主界面入口", ""),
        ("│  (实现Calculator)    │", ""),
        ("└─────────┬───────────┘", ""),
        ("          │ creates", ""),
        ("          ▼", ""),
        ("┌─────────────────────┐", ""),
        ("│   CalculatorImpl    │ ◄── 核心计算引擎", ""),
        ("│  (~20个方法)         │", ""),
        ("└──┬──────┬───────────┘", ""),
        ("   │      │ uses", ""),
        ("   ▼      ▼", ""),
        ("┌──────┐ ┌────────────┐", ""),
        ("│EvalEx│ │NumberFormat│", ""),
        ("│求值  │ │Helper格式化│", ""),
        ("└──────┘ └────────────┘", ""),
        ("", ""),
        ("┌─────────────────────┐", ""),
        ("│  Converter 接口     │ ◄── 策略模式基类", ""),
        ("└─────────┬───────────┘", ""),
        ("  implemented by (9种)", ""),
        ("  Length / Area / Volume / Mass / Temperature", ""),
        ("  Time / Speed / Pressure / Energy", ""),
        ("", ""),
        ("┌─────────────────────┐", ""),
        ("│ CalculatorDatabase  │ ◄── Room单例", ""),
        ("│  └─ CalculatorDao   │", ""),
        ("│     CRUD 操作       │", ""),
        ("└─────────┬───────────┘", ""),
        ("          │ manages", ""),
        ("          ▼", ""),
        ("┌─────────────────────┐", ""),
        ("│   History Entity    │", ""),
        ("│ (formula/result/ts) │", ""),
        ("└─────────────────────┘", ""),
    ]
    y = 1.16
    for line, _ in class_relations:
        add_text_box(slide, Inches(0.38), Inches(y), Inches(5.3), Inches(0.165),
                     line, 8.5, False, DARK_GRAY, PP_ALIGN.LEFT, "Consolas")
        y += 0.155
    
    # 右侧：设计模式识别表
    pattern_data = [
        ["模式名称", "应用位置", "说明"],
        ["MVP 变体", "MainActivity + CalculatorImpl", "View-Presenter 分离"],
        ["策略模式 Strategy", "Converter 接口 + 9种实现", "单位转换算法封装"],
        ["单例模式 Singleton", "CalculatorDatabase", "全局唯一数据库实例"],
        ["观察者模式 Observer", "OnUnitChangedListener", "单位变更通知回调"],
        ["工厂模式 Factory", "ValueWithUnit 创建", "统一创建带单位的值对象"],
    ]
    add_table(slide, len(pattern_data), 3, Inches(5.92), Inches(0.78), Inches(3.98), Inches(1.9), pattern_data,
              col_widths=[Inches(1.4), Inches(1.6), Inches(0.98)], header_font_size=10, cell_font_size=9)
    
    # 类图颜色说明
    legend_data = [
        ["颜色标识", "含义"],
        ["浅绿色 #E8F5E9", "接口 (Interface)"],
        ["浅蓝色 #E3F2FD", "Activity / UI层"],
        ["浅橙色 #FFF3E0", "核心业务逻辑"],
        ["浅紫色 #F3E5F5", "数据层 (DB/Entity)"],
    ]
    add_table(slide, len(legend_data), 2, Inches(5.92), Inches(2.8), Inches(3.98), Inches(1.4), legend_data,
              col_widths=[Inches(1.6), Inches(2.38)], header_font_size=10, cell_font_size=9)


def create_function_modules(prs):
    """P05 功能模块"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    add_title_shape(slide, "四、系统功能结构图", 0, 0, prs.slide_width, Inches(0.65))
    
    # 功能模块概览（左侧）
    modules = [
        ("🧮 计算器模块", ["基础运算: 加减乘除、小数、负数", "高级运算: 幂运算^、开方√、百分比%", "数字格式化: 千分位分隔、本地化"]),
        ("🔄 单位转换模块", ["长度转换 17种 / 面积转换 10种", "体积转换 24种 / 质量转换 15种", "温度/时间/速度/压强/能量 共48种"]),
        ("📊 历史记录模块", ["保存计算记录、查看历史列表", "复用历史结果、清空历史记录"]),
        ("⚙️ 设置模块", ["振动开关、屏幕常亮", "主题颜色定制、语言切换"]),
        ("🖼️ 小部件模块", ["桌面快捷计算、自定义配色"]),
    ]
    
    x_pos = 0.25
    for title, items in modules:
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     Inches(x_pos), Inches(0.82), Inches(1.82), Inches(2.5))
        set_shape_fill(box, RGBColor(0xE3, 0xF2, 0xFD))
        add_text_box(slide, Inches(x_pos + 0.08), Inches(0.88), Inches(1.66), Inches(0.28),
                     title, 11, True, ACCENT_BLUE, PP_ALIGN.LEFT)
        y = 1.2
        for item in items:
            add_text_box(slide, Inches(x_pos + 0.08), Inches(y), Inches(1.66), Inches(0.35),
                         f"• {item}", 8.5, False, DARK_GRAY, PP_ALIGN.LEFT)
            y += 0.36
        x_pos += 1.9
    
    # 数据流示意（下方）
    flow_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      Inches(0.2), Inches(3.45), Inches(9.4), Inches(2.3))
    set_shape_fill(flow_box, RGBColor(0xFC, 0xFC, 0xFC))
    
    add_text_box(slide, Inches(0.35), Inches(3.53), Inches(4), Inches(0.28),
                 "核心数据流", 12, True, ACCENT_BLUE, PP_ALIGN.LEFT)
    
    flow_text = (
        "用户点击按钮\n"
        "    ↓\n"
        "MainActivity.numpadClicked()\n"
        "    ↓\n"
        "CalculatorImpl 构建表达式\n"
        "    ↓\n"
        "EvalEx 解析表达式 → BigDecimal 结果\n"
        "    ↓\n"
        "NumberFormatHelper 格式化显示\n"
        "    ↓\n"
        "showNewResult() 更新 UI\n\n"
        "同时: HistoryHelper → Room Database 持久化历史记录"
    )
    add_text_box(slide, Inches(0.4), Inches(3.85), Inches(4.3), Inches(1.8),
                 flow_text, 9.5, False, DARK_GRAY, PP_ALIGN.LEFT, "Consolas")
    
    # 单位转换详细清单（右侧下）
    conv_data = [
        ["转换器类别", "单位数", "基准单位", "代表单位"],
        ["LengthConverter", "17 种", "Meter", "km/cm/mm/inch/mile/光年..."],
        ["AreaConverter", "10 种", "SquareMeter", "km²/inch²/公顷/亩"],
        ["VolumeConverter", "24 种", "CubicMeter", "L/mL/gal/桶/盎司"],
        ["MassConverter", "15 种", "Kilogram", "g/lb/oz/t/原子质量"],
        ["TemperatureConverter", "4 种", "特殊偏移", "℃ / ℉ / K / °R"],
        ["TimeConverter", "7 种", "Second", "h/min/ms/d/周/y"],
        ["SpeedConverter", "8 种", "m/s", "km/h/马赫/光速/节"],
        ["PressureConverter", "10 种", "Pascal", "bar/atm/psi/托"],
        ["EnergyConverter", "14 种", "Joule", "cal/kWh/eV/BTU"],
        ["合计", "109 种", "—", "覆盖物理/工程/日常生活"],
    ]
    add_table(slide, len(conv_data), 4, Inches(4.85), Inches(3.52), Inches(4.75), Inches(2.2), conv_data,
              col_widths=[Inches(1.35), Inches(0.6), Inches(0.95), Inches(1.85)],
              header_font_size=9, cell_font_size=8)


def create_test_system(prs):
    """P06 测试体系设计"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    add_title_shape(slide, "五、测试体系设计与架构", 0, 0, prs.slide_width, Inches(0.65))
    
    # 测试金字塔（左侧）- 用形状绘制
    pyramid_colors = [
        (RGBColor(0xC8, 0xE6, 0xC9), "静态分析 / Lint (Detekt/Jacoco/Lint)", "底层：速度快、数量多、成本低"),
        (RGBColor(0xFF, 0xEC, 0xB3), "单元测试 (JUnit + MockK + Robolectric)", "~70+ 用例 | JVM 本地运行 | 秒级"),
        (RGBColor(0xFF, 0xCC, 0xBC), "集成测试 (Espresso + AndroidX Test)", "~21 用例 | 真机/模拟器 | 分钟级"),
        (RGBColor(0xFF, 0xCD, 0xD2), "端到端测试 E2E (少量)", "顶层：速度慢、数量少、成本高"),
    ]
    
    base_left = 0.35
    base_width_step = 1.6
    pyramid_top = 5.3
    for i, (color, label, detail) in enumerate(pyramid_colors):
        width = 1.8 + i * base_width_step
        left = base_left + ((4.5 - width) / 2)
        height = 0.85
        top = pyramid_top - (len(pyramid_colors) - i) * height
        
        shape = slide.shapes.add_shape(MSO_SHAPE.TRAPEZOID, Inches(left), Inches(top),
                                       Inches(width), Inches(height))
        set_shape_fill(shape, color)
        
        add_text_box(slide, Inches(left + 0.1), Inches(top + 0.08), Inches(width - 0.2), Inches(0.35),
                     label, 9, True, BLACK, PP_ALIGN.CENTER)
        add_text_box(slide, Inches(left + 0.1), Inches(top + 0.42), Inches(width - 0.2), Inches(0.3),
                     detail, 7.5, False, DARK_GRAY, PP_ALIGN.CENTER)
    
    # 右侧上方：框架选型
    framework_data = [
        ["层级", "技术栈"],
        ["单元测试", "JUnit 4 + MockK + Robolectric"],
        ["集成测试", "Espresso + AndroidX Test (ActivityScenarioRule)"],
        ["静态分析", "Detekt Kotlin Lint + Jacoco + Android Lint"],
        ["覆盖率", "Jacoco (单元+集成合并报告)"],
    ]
    add_table(slide, len(framework_data), 2, Inches(5.2), Inches(0.78), Inches(4.45), Inches(1.85), framework_data,
              col_widths=[Inches(1.2), Inches(3.25)], header_font_size=10, cell_font_size=9)
    
    # 右侧中间：环境配置
    env_data = [
        ["配置项", "值", "说明"],
        ["编译 SDK", "compileSdk 34", "编译时使用的 SDK 版本"],
        ["目标 SDK", "targetSdk 34", "目标设备版本"],
        ["最低 SDK", "minSdk 26", "最低兼容 Android 8.0"],
        ["Robolectric SDK", "@Config(sdk=[33])", "模拟 SDK 版本"],
        ["Kotlin 版本", "1.9.22", "编译器版本"],
        ["Gradle / AGP", "8.4 / 8.2.0", "构建工具"],
    ]
    add_table(slide, len(env_data), 3, Inches(5.2), Inches(2.73), Inches(4.45), Inches(2.0), env_data,
              col_widths=[Inches(1.3), Inches(1.3), Inches(1.85)], header_font_size=9, cell_font_size=8)
    
    # 底部：测试分类策略
    classify_data = [
        ["测试类型", "执行位置", "运行环境", "耗时", "覆盖重点"],
        ["纯逻辑单元测试", "test/", "JVM 本地", "秒级", "纯 Kotlin 逻辑"],
        ["Android 组件测试", "test/", "Robolectric 模拟", "秒~十秒级", "Activity/View/Context"],
        ["数据库测试", "test/", "Robolectric + Room 内存", "十秒级", "DAO CRUD / 单例线程安全"],
        ["UI 集成测试", "androidTest/", "真机/模拟器", "分钟级", "用户交互 / E2E 场景"],
    ]
    add_table(slide, len(classify_data), 5, Inches(0.25), Inches(4.88), Inches(9.4), Inches(1.72), classify_data,
              col_widths=[Inches(1.6), Inches(1.1), Inches(1.8), Inches(1.0), Inches(3.9)],
              header_font_size=10, cell_font_size=9)


def create_unit_tests(prs):
    """P07 单元测试设计"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    add_title_shape(slide, "六、单元测试设计与典型案例", 0, 0, prs.slide_width, Inches(0.65))
    
    # 顶部统计总表
    stat_data = [
        ["测试模块", "测试类", "方法数", "覆盖重点", "测试框架"],
        ["activities", "1", "7", "MainActivity UI 交互", "Robolectric"],
        ["helpers 核心", "3", "~15", "CalculatorImpl / Constants / Format", "MockK + JUnit"],
        ["helpers.converters", "9", "~25", "全部 9 种转换器", "纯 JUnit (JVM)"],
        ["databases", "1", "6", "单例 & DAO 操作", "Robolectric + Room"],
        ["adapters/dialogs/models", "4", "~19", "Adapter/Dialog/State", "JUnit / Mockito"],
        ["合计", "**18**", "**~70+**", "全部核心模块", "—"],
    ]
    add_table(slide, len(stat_data), 5, Inches(0.2), Inches(0.75), Inches(9.4), Inches(1.85), stat_data,
              col_widths=[Inches(1.7), Inches(0.7), Inches(0.8), Inches(3.2), Inches(3.0)],
              header_font_size=10, cell_font_size=9)
    
    # 四个典型案例卡片
    cases = [
        ("案例一: MainActivity Robolectric", GREEN,
         ["使用 @RunWith(RobolectricTestRunner::class)",
          "@Config(sdk=[33]) + @LooperMode(PAUSED)",
          "6个方法: 加减乘除 + 清除操作",
          "shadowOf(Looper).idle() 确保 UI 更新完成",
          "验证: assertEquals('3', result)"]),
        ("案例二: CalculatorImpl MockK", PURPLE,
         ["mockk(relaxed=true) 宽松 Mock",
          "mockkConstructor(NumberFormatHelper)",
          "拦截 formatForDisplay 返回原值",
          "verify(exactly=N) 精确控制调用次数",
          "隔离 NumberFormatHelper 外部依赖"]),
        ("案例三: 数据库单例 & DAO", ORANGE,
         ["单例唯一性: assertSame(inst1, inst2)",
          "线程安全: CountDownLatch 10线程并发",
          "销毁重建: destroyInstance 后重新获取",
          "DAO CRUD: 插入/查询/LIMIT/批量/清空",
          "Room 内存数据库独立测试文件"]),
        ("案例四: Converter Roundtrip", TEAL,
         ["纯 JVM 测试，无需 Android 依赖",
          "Roundtrip: value → toBase → fromBase == value",
          "排除 TemperatureConverter (非线性)",
          "零值/同单位恒等性验证",
          "BigDecimal compareTo==0 精确断言"]),
    ]
    
    card_width = 4.55
    card_height = 2.0
    positions = [(0.23, 2.72), (4.93, 2.72), (0.23, 4.83), (4.93, 4.83)]
    
    for idx, ((title, color, items), (lx, ty)) in enumerate(zip(cases, positions)):
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     Inches(lx), Inches(ty), Inches(card_width), Inches(card_height))
        set_shape_fill(box, RGBColor(0xFA, 0xFA, 0xFA))
        box.line.color.rgb = color
        
        # 标题条
        title_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                           Inches(lx), Inches(ty), Inches(card_width), Inches(0.33))
        set_shape_fill(title_bar, color)
        add_text_box(slide, Inches(lx + 0.08), Inches(ty + 0.03), Inches(card_width - 0.16), Inches(0.28),
                     title, 10, True, WHITE, PP_ALIGN.LEFT)
        
        # 内容要点
        y = ty + 0.4
        for item in items:
            add_text_box(slide, Inches(lx + 0.1), Inches(y), Inches(card_width - 0.2), Inches(0.27),
                         f"✓ {item}", 8.5, False, DARK_GRAY, PP_ALIGN.LEFT)
            y += 0.29


def create_integration_tests(prs):
    """P08 集成测试设计"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    add_title_shape(slide, "七、Espresso UI 集成测试设计", 0, 0, prs.slide_width, Inches(0.65))
    
    # 左侧：集成测试统计
    integ_stat = [
        ["测试类名", "方法数", "场景", "页面"],
        ["MainActivityEspressoTest", "7", "计算、清除、长按清空", "主计算器界面"],
        ["SettingsActivityEspressoTest", "~3", "设置项切换", "设置页面"],
        ["UnitConverterActivityEspressoTest", "5", "输入、清除、转换显示", "单位转换页面"],
        ["UnitConverterPickerActivityTest", "~2", "类型选择", "单位选择器"],
        ["WidgetConfigureActivityTest", "~4", "颜色配置", "小部件配置页"],
        ["合计", "**~21**", "5 个 Activity", "通过率 ~90%+"],
    ]
    add_table(slide, len(integ_stat), 4, Inches(0.2), Inches(0.78), Inches(4.5), Inches(2.35), integ_stat,
              col_widths=[Inches(2.0), Inches(0.6), Inches(1.2), Inches(0.7)],
              header_font_size=10, cell_font_size=9)
    
    # Espresso vs Robolectric 对比（右侧上）
    compare_data = [
        ["对比维度", "Robolectric 单元测试", "Espresso 集成测试"],
        ["运行环境", "JVM (模拟 Android)", "真机 / 模拟器"],
        ["UI 交互", "view.performClick() 直接调用", "onView(...).perform(click())"],
        ["断言方式", "assertEquals(text, result)", "check(matches(withText()))"],
        ["生命周期", "手动 create/start/resume", "ActivityScenarioRule 自动管理"],
        ["执行速度", "~100ms / 测试", "~1-3s / 测试"],
        ["适用场景", "逻辑验证 / CI 快速反馈", "真实交互 / E2E 验收"],
    ]
    add_table(slide, len(compare_data), 3, Inches(4.82), Inches(0.78), Inches(4.8), Inches(2.35), compare_data,
              col_widths=[Inches(1.3), Inches(1.75), Inches(1.75)],
              header_font_size=9, cell_font_size=8.5)
    
    # 典型代码示例（左侧下）
    code_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      Inches(0.2), Inches(3.23), Inches(4.5), Inches(2.55))
    set_shape_fill(code_box, RGBColor(0x1E, 0x1E, 0x1E))
    
    code_text = (
        "// MainActivityEspressoTest.kt\n"
        "@RunWith(AndroidJUnit4::class)\n"
        "class MainActivityEspressoTest {\n"
        "    @get:Rule\n"
        "    val activityRule = ActivityScenarioRule(\n"
        "        MainActivity::class.java)\n\n"
        "    @Test\n"
        "    fun testAddTwoNumbers() {\n"
        "        clickView(R.id.btn_1)\n"
        "        clickView(R.id.btn_plus)\n"
        "        clickView(R.id.btn_2)\n"
        "        clickView(R.id.btn_equals)\n\n"
        "        onView(withId(R.id.result))\n"
        "            .check(matches(withText(\"3\")))\n"
        "    }\n"
        "\n"
        "    @Test fun testClearAllOnLongPress() { ... }\n"
        "}"
    )
    add_text_box(slide, Inches(0.32), Inches(3.31), Inches(4.26), Inches(2.4),
                 code_text, 8, False, RGBColor(0xD4, 0xD4, 0xD4), PP_ALIGN.LEFT, "Consolas")
    
    # 执行流程（右侧下）
    flow_data = [
        ["步骤", "操作"],
        ["① 安装", "安装 APK 到测试设备"],
        ["② 启动", "ActivityScenarioRule 自动启动 Activity"],
        ["③ 定位", "Espresso onView(withId(...)) 定位控件"],
        ["④ 操作", "perform(click() / longClick()) 模拟用户行为"],
        ["⑤ 断言", "check(matches(withText(...))) 验证结果"],
        ["⑥ 判定", "PASS → 下一个 | FAIL → 截图+记录"],
        ["⑦ 报告", "生成 HTML/XML 测试报告"],
    ]
    add_table(slide, len(flow_data), 2, Inches(4.82), Inches(3.23), Inches(4.8), Inches(2.55), flow_data,
              col_widths=[Inches(0.8), Inches(4.0)], header_font_size=10, cell_font_size=9)


def create_coverage_requirements(prs):
    """P09 覆盖率要求"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    add_title_shape(slide, "八、测试要求与覆盖率标准", 0, 0, prs.slide_width, Inches(0.65))
    
    # 左侧：测试要求规范
    req_data = [
        ["要求类别", "指标", "目标值", "状态"],
        ["功能正确性", "核心业务逻辑", "100% 正确", "✅ 已达成"],
        ["Line Coverage", "行覆盖率", "≥ 80%", "⚠️ 待提升"],
        ["Method Coverage", "方法覆盖率", "≥ 80%", "⚠️ 待提升"],
        ["Branch Coverage", "分支覆盖率", "≥ 70%", "❌ 未达标"],
        ["Class Coverage", "类覆盖率", "≥ 90%", "⚠️ 待提升"],
        ["边界条件", "边界值测试", "全覆盖", "⚠️ 部分"],
        ["异常处理", "异常路径", "全覆盖", "⚠️ 部分"],
    ]
    add_table(slide, len(req_data), 4, Inches(0.2), Inches(0.78), Inches(4.5), Inches(2.9), req_data,
              col_widths=[Inches(1.4), Inches(1.3), Inches(1.0), Inches(0.8)],
              header_font_size=10, cell_font_size=9)
    
    # 右侧上方：质量金字塔（文本表示）
    quality_box = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE,
                                         Inches(6.5), Inches(0.85), Inches(2.8), Inches(2.5))
    # 使用矩形堆叠替代三角形
    q_levels = [
        (Inches(5.5), Inches(0.85), Inches(3.8), Inches(0.4), RGBColor(0xFF, 0xCD, 0xD2), "可靠性 — 无间歇性失败"),
        (Inches(5.7), Inches(1.28), Inches(3.4), Inches(0.4), RGBColor(0xFF, 0xEC, 0xB3), "可重复性 — 任意次执行一致"),
        (Inches(5.9), Inches(1.71), Inches(3.0), Inches(0.4), RGBColor(0xFF, 0xE0, 0xB2), "独立性 — 测试间无依赖"),
        (Inches(6.1), Inches(2.14), Inches(2.6), Inches(0.4), RGBColor(0xFE, 0xE0, 0xB2), "快速执行 — 秒级完成"),
        (Inches(6.3), Inches(2.57), Inches(2.2), Inches(0.4), RGBColor(0xC8, 0xE6, 0xC9), "可读性强 — 命名清晰"),
        (Inches(6.5), Inches(3.0), Inches(1.8), Inches(0.4), RGBColor(0xB3, 0xE5, 0xFC), "全面性与针对性平衡"),
    ]
    for lx, ly, lw, lh, lc, lt in q_levels:
        s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, lx, ly, lw, lh)
        set_shape_fill(s, lc)
        add_text_box(slide, lx + Inches(0.05), ly + Inches(0.07), lw - Inches(0.1), lh - Inches(0.05),
                     lt, 8.5, False, BLACK, PP_ALIGN.CENTER)
    
    # Jacoco 配置（下方）
    jacoco_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                        Inches(0.2), Inches(3.8), Inches(9.4), Inches(1.95))
    set_shape_fill(jacoco_box, RGBColor(0x2D, 0x2D, 0x2D))
    
    add_text_box(slide, Inches(0.35), Inches(3.88), Inches(4), Inches(0.28),
                 "Jacoco 配置 (build.gradle.kts)", 11, True, RGBColor(0x4A, 0xF5, 0x9B), PP_ALIGN.LEFT)
    
    jacoco_code = (
        "android {\n"
        "    buildTypes {\n"
        "        debug {\n"
        "            enableUnitTestCoverage = true\n"
        "            enableAndroidTestCoverage = true\n"
        "        }\n"
        "    }\n"
        "}\n\n"
        "# 排除规则:\n"
        "#   • 自动生成代码 (BuildConfig, R类, DataBinding)\n"
        "#   • 接口默认方法\n"
        "#   • 仅含 getter/setter 的简单数据类\n\n"
        "# 生成命令:\n"
        "#   ./gradlew testDebugUnitTestJacocoReport\n"
        "#   ./gradlew createDebugCoverageReport"
    )
    add_text_box(slide, Inches(0.35), Inches(4.2), Inches(9.1), Inches(1.5),
                 jacoco_code, 9, False, RGBColor(0xDD, 0xDD, 0xDD), PP_ALIGN.LEFT, "Consolas")


def create_coverage_analysis(prs):
    """P10 覆盖率统计分析"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    add_title_shape(slide, "九、覆盖率统计与模块分析", 0, 0, prs.slide_width, Inches(0.65))
    
    # 总体覆盖率仪表盘（顶部）
    metrics = [
        ("Class", "78.4%", "(131/167)", RGBColor(0xFB, 0x8C, 0x00)),
        ("Method", "42.2%", "(191/453)", RED),
        ("Branch", "0%", "(0/598)", RED),
        ("Line", "35.3%", "(746/2116)", RED),
    ]
    mx = 0.35
    for label, value, detail, color in metrics:
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     Inches(mx), Inches(0.78), Inches(2.25), Inches(1.05))
        set_shape_fill(box, RGBColor(0xFA, 0xFA, 0xFA))
        box.line.color.rgb = color
        box.line.width = Pt(2)
        add_text_box(slide, Inches(mx + 0.05), Inches(0.84), Inches(2.15), Inches(0.25),
                     label, 11, True, DARK_GRAY, PP_ALIGN.CENTER)
        add_text_box(slide, Inches(mx + 0.05), Inches(1.08), Inches(2.15), Inches(0.4),
                     value, 26, True, color, PP_ALIGN.CENTER)
        add_text_box(slide, Inches(mx + 0.05), Inches(1.5), Inches(2.15), Inches(0.25),
                     detail, 9, False, DARK_GRAY, PP_ALIGN.CENTER)
        mx += 2.35
    
    # 模块覆盖率明细表（左下）
    module_data = [
        ["Package", "Class%", "Method%", "Branch%", "Line%", "状态"],
        ["helpers.converters", "99.2", "96.4", "N/A", "98.2", "✅ 优秀"],
        ["helpers", "12.5", "1.8", "0%", "0.2", "⚠️ 待补充"],
        ["activities", "0", "0", "0%", "0", "❌ 未覆盖"],
        ["views", "0", "0", "0%", "0", "❌ 未覆盖"],
        ["databases/compose/adapters", "0", "0", "0%", "0", "❌ 未覆盖"],
        ["其他 (6个包)", "0", "0", "0%", "0", "❌ 未覆盖"],
    ]
    add_table(slide, len(module_data), 6, Inches(0.2), Inches(1.95), Inches(5.5), Inches(2.35), module_data,
              col_widths=[Inches(1.5), Inches(0.7), Inches(0.75), Inches(0.7), Inches(0.7), Inches(1.15)],
              header_font_size=9, cell_font_size=8.5)
    
    # 原因分析 + 改进计划（右侧）
    analysis_items = [
        "为什么 converters 最高？",
        "  ✓ 无 Android 依赖，纯 JVM 逻辑",
        "  ✓ 策略模式优势，独立易测",
        "  ✓ Roundtrip 一个方法覆盖全部单位",
        "  ✓ BigDecimal 精确比对避免浮点误差",
        "",
        "为什么其他模块低？",
        "  ✗ activities/views 需要 Robolectric/Espresso",
        "  ✗ databases 已有测试但未纳入 Jacoco",
        "  ✗ compose 需要 Compose Testing",
        "",
        "改进计划:",
        "  P0: 完善 activities Robolectric → Line +30%",
        "  P0: databases DAO 纳入 Jacoco → Line +5%",
        "  P1: views ConverterView 测试 → Line +10%",
        "  P1: helpers 增强 CalculatorImpl → Line +8%",
    ]
    add_bullet_points(slide, Inches(5.82), Inches(1.95), Inches(3.85), Inches(2.35),
                      analysis_items, 8.5, DARK_GRAY)
    
    # 覆盖率柱状图（底部）- 用文本模拟
    chart_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                       Inches(0.2), Inches(4.4), Inches(9.4), Inches(1.4))
    set_shape_fill(chart_box, RGBColor(0xF8, 0xF9, 0xFA))
    
    add_text_box(slide, Inches(0.35), Inches(4.47), Inches(3), Inches(0.25),
                 "各模块行覆盖率对比", 10, True, ACCENT_BLUE, PP_ALIGN.LEFT)
    
    bars = [
        ("helpers.converters", 98.2, GREEN),
        ("helpers", 0.2, RED),
        ("activities", 0, RED),
        ("views", 0, RED),
        ("databases", 0, RED),
        ("compose", 0, RED),
        ("others", 0, RED),
    ]
    bx = 0.4
    max_bar_width = 4.0
    for name, pct, color in bars:
        bar_w = max(pct / 100 * max_bar_width, 0.02)
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                     Inches(bx), Inches(4.8), Inches(bar_w), Inches(0.28))
        set_shape_fill(bar, color)
        add_text_box(slide, Inches(bx), Inches(5.1), Inches(max_bar_width + 0.3), Inches(0.2),
                     f"{name}: {pct}%", 7.5, False, DARK_GRAY, PP_ALIGN.LEFT)
        bx += 1.3


def create_performance_test(prs):
    """P11 性能测试"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    add_title_shape(slide, "十、性能测试与瓶颈分析", 0, 0, prs.slide_width, Inches(0.65))
    
    # 性能基准数据表（左侧）
    perf_data = [
        ["性能指标", "基准值", "实测值", "状态"],
        ["冷启动时间", "< 500ms", "~380ms", "✅ 优秀"],
        ["热启动时间", "< 200ms", "~120ms", "✅ 优秀"],
        ["常规计算响应", "< 16ms (1帧)", "< 10ms", "✅ 优秀"],
        ["复杂计算响应", "< 100ms", "~45ms", "✅ 良好"],
        ["FPS (交互状态)", ">= 55fps", "58-60fps", "✅ 流畅"],
        ["内存占用 (空闲)", "< 50MB", "~35MB", "✅ 正常"],
        ["内存占用 (峰值)", "< 150MB", "~85MB", "✅ 良好"],
        ["CPU 占用 (密集)", "< 80%", "~45%", "✅ 正常"],
        ["GC 频率", "< 5次/min", "~2-3次", "✅ 正常"],
    ]
    add_table(slide, len(perf_data), 4, Inches(0.2), Inches(0.78), Inches(4.5), Inches(3.3), perf_data,
              col_widths=[Inches(1.4), Inches(1.0), Inches(1.0), Inches(1.1)],
              header_font_size=9, cell_font_size=8.5)
    
    # CPU 热点分布（右侧上）
    hotspots = [
        ("CalculatorImpl.numpadClicked()", "28%", "  ├ StringBuilder.append 8%"),
        ("EvalEx Expression.evaluate()", "22%", "  ├ 词法 Lexer 10% / Parser 7%"),
        ("ConverterView.convertValue()", "18%", "  ├ toBase 10% / fromBase 8%"),
        ("Room DB Operation", "8%", "  ├ INSERT/UPDATE 5% / SELECT 3%"),
        ("UI Rendering (Measure/Layout)", "9%", ""),
        ("Other (GC/System)", "15%", ""),
    ]
    
    hot_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      Inches(4.82), Inches(0.78), Inches(4.68), Inches(2.6))
    set_shape_fill(hot_box, RGBColor(0xFD, 0xF8, 0xF0))
    add_text_box(slide, Inches(4.97), Inches(0.86), Inches(3), Inches(0.25),
                 "CPU 热点分布", 11, True, ORANGE, PP_ALIGN.LEFT)
    
    hy = 1.15
    for name, pct, sub in hotspots:
        add_text_box(slide, Inches(4.97), Inches(hy), Inches(2.8), Inches(0.22),
                     name, 8.5, False, DARK_GRAY, PP_ALIGN.LEFT)
        add_text_box(slide, Inches(7.8), Inches(hy), Inches(0.6), Inches(0.22),
                     pct, 9, True, RED, PP_ALIGN.RIGHT)
        if sub:
            add_text_box(slide, Inches(4.97), Inches(hy + 0.2), Inches(4.3), Inches(0.2),
                         sub, 7.5, False, RGBColor(0x99, 0x99, 0x99), PP_ALIGN.LEFT)
        hy += 0.37
    
    # 优化建议表（右侧下）
    opt_data = [
        ["瓶颈位置", "严重程度", "优化建议"],
        ["UI 线程阻塞", "🔴 高", "将 EvalEx 求值移至协程/后台线程"],
        ["NumberFormatHelper", "🟡 中", "引入缓存机制，仅在值变化时格式化"],
        ["ConverterView 实时转换", "🟡 中", "加入防抖 debounce 300ms 延迟后转换"],
        ["Room 写入频繁", "🟡 中", "批量写入或 Write-Ahead Logging"],
        ["EvalEx 解析开销", "🟢 低", "固定模式预编译表达式模板"],
    ]
    add_table(slide, len(opt_data), 3, Inches(4.82), Inches(3.5), Inches(4.68), Inches(1.95), opt_data,
              col_widths=[Inches(1.5), Inches(0.85), Inches(2.33)],
              header_font_size=9, cell_font_size=8)
    
    # 测试场景（下方）
    scenario_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                          Inches(0.2), Inches(4.2), Inches(4.5), Inches(1.58))
    set_shape_fill(scenario_box, RGBColor(0xE8, 0xAF, 0xCA))
    
    scen_data = [
        ["场景编号", "名称", "关注指标"],
        ["S01", "冷启动测试 (<500ms)", "启动耗时"],
        ["S02", "连续基础计算 ×100次", "CPU / 响应速度"],
        ["S03", "连续复杂计算 ×100次", "CPU / 内存"],
        ["S04", "大数精度计算 (999M×999M)", "内存 / 时间"],
        ["S05", "单位转换压力测试", "UI 响应 / 抖动"],
        ["S06", "历史记录插入 ×1000条", "DB IO / 内存"],
        ["S07", "长时间运行稳定性 10min", "泄漏 / GC"],
    ]
    add_table(slide, len(scen_data), 3, Inches(0.25), Inches(4.28), Inches(4.4), Inches(1.48), scen_data,
              col_widths=[Inches(0.6), Inches(2.4), Inches(1.4)],
              header_font_size=8.5, cell_font_size=7.5)
    
    # 优化路线图（右下）
    roadmap_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                         Inches(4.82), Inches(5.55), Inches(4.68), Inches(0.35))
    set_shape_fill(roadmap_box, RGBColor(0xDC, 0xED, 0xC8))
    add_text_box(slide, Inches(4.97), Inches(5.59), Inches(4.38), Inches(0.28),
                 "路线图: Phase1 协程化+防抖 → Phase2 缓存+WAL → Phase3 Compose迁移",
                 8.5, False, GREEN, PP_ALIGN.CENTER)


def create_summary(prs):
    """P12 总结与展望"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    add_title_shape(slide, "十一、测试总结与未来展望", 0, 0, prs.slide_width, Inches(0.65))
    
    # 成果总览仪表盘（顶部）
    results = [
        ("测试规模", "18类 / 91+用例", GREEN),
        ("整体通过率", "~95%+", GREEN),
        ("行覆盖率", "35.3%", ORANGE),
        ("冷启动", "~380ms", GREEN),
    ]
    rx = 0.35
    for label, value, color in results:
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     Inches(rx), Inches(0.78), Inches(2.25), Inches(0.85))
        set_shape_fill(box, color)
        add_text_box(slide, Inches(rx + 0.05), Inches(0.83), Inches(2.15), Inches(0.25),
                     label, 10, True, WHITE, PP_ALIGN.CENTER)
        add_text_box(slide, Inches(rx + 0.05), Inches(1.1), Inches(2.15), Inches(0.42),
                     value, 20, True, WHITE, PP_ALIGN.CENTER)
        rx += 2.35
    
    # 覆盖矩阵（左侧）
    matrix_data = [
        ["功能模块", "单元测试", "集成测试", "覆盖状态"],
        ["基础四则运算", "✅", "✅", "✅ 完整"],
        ["数字格式化 (i18n)", "✅", "✅", "✅ 完整"],
        ["单位转换 (9类)", "✅ 完整", "✅", "✅ **优秀 98.2%**"],
        ["计算历史 CRUD", "✅", "⚠️ 部分", "✅ 基本完整"],
        ["UI 交互 (按钮/清除)", "✅", "✅", "✅ 完整"],
        ["高级运算 (幂/根/%)", "⚠️ 部分", "❌", "⚠️ 待完善"],
        ["设置页面 (Compose)", "❌", "⚠️ 部分", "⚠️ 基础"],
        ["小部件 Widget", "❌", "⚠️ 部分", "⚠️ 基础"],
    ]
    add_table(slide, len(matrix_data), 4, Inches(0.2), Inches(1.75), Inches(4.5), Inches(3.0), matrix_data,
              col_widths=[Inches(1.6), Inches(0.9), Inches(0.9), Inches(1.1)],
              header_font_size=9, cell_font_size=8.5)
    
    # 挑战与解决方案（中右）
    challenge_data = [
        ["挑战", "解决方案", "经验教训"],
        ["Android 依赖隔离困难", "MockK mock 构造函数", "精心设计 Mock 策略"],
        ["Robolectric vs 真机差异", "双版本互补策略", "单元验证逻辑 / 集成验证行为"],
        ["Jacoco 覆盖率合并问题", "分别报告交叉验证", "排除自动生成代码"],
        ["单例测试污染", "独立 DB 文件 + tearDown", "测试隔离至关重要"],
        ["温度非线性转换", "过滤 TemperatureConverter", "特殊边界需单独处理"],
        ["UI 测试不稳定", "IdleResource 显式等待", "避免硬编码 sleep"],
    ]
    add_table(slide, len(challenge_data), 3, Inches(4.82), Inches(1.75), Inches(4.68), Inches(2.4), challenge_data,
              col_widths=[Inches(1.5), Inches(1.7), Inches(1.48)],
              header_font_size=9, cell_font_size=8)
    
    # 最终结论卡片（下方大框）
    conclusion_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                            Inches(0.2), Inches(4.85), Inches(9.4), Inches(1.0))
    set_shape_fill(conclusion_box, MID_BLUE)
    
    conclusion_lines = [
        ("✅ 功能正确性:", "核心四则运算、单位转换、CRUD 操作均符合预期"),
        ("⚠️ 测试覆盖率:", "converters 模块 98.2% 优异；其他模块待提升 (整体 35.3%)"),
        ("✅ 性能表现:", "冷启动 ~380ms，内存 35-85MB，CPU 负载健康，优于行业基准"),
        ("📊 总评:", "系统功能基本正确，测试框架搭建完成，具备生产就绪潜力"),
    ]
    cx = 0.38
    for label, text in conclusion_lines:
        color = RGBColor(0xBB, 0xDE, 0xFB) if "✅" in label or "📊" in label else RGBColor(0xFF, 0xEC, 0xB3)
        add_text_box(slide, Inches(cx), Inches(4.93), Inches(1.3), Inches(0.22),
                     label, 9, True, color, PP_ALIGN.LEFT)
        add_text_box(slide, Inches(cx + 1.35), Inches(4.93), Inches(7.8), Inches(0.22),
                     text, 9, False, WHITE, PP_ALIGN.LEFT)
        cx = 0.38
        if label != conclusion_lines[-1][0]:
            pass  # keep on same line area, use different approach
        # Actually let's stack them vertically
    # Redraw conclusions stacked
    cy = 4.93
    for label, text in conclusion_lines:
        lcolor = RGBColor(0xBB, 0xDE, 0xFB) if "✅" in label or "📊" in label else RGBColor(0xFF, 0xEC, 0xB3)
        add_text_box(slide, Inches(0.38), Inches(cy), Inches(9), Inches(0.2),
                     f"{label}  {text}", 9, False, WHITE, PP_ALIGN.LEFT)
        cy += 0.185


# ======================== 主程序 ========================

def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)  # 16:10 宽屏比例
    
    print("正在生成 PPT...")
    print("  [1/12] 封面...")
    create_cover(prs)
    print("  [2/12] 项目背景...")
    create_background(prs)
    print("  [3/12] 系统架构...")
    create_architecture(prs)
    print("  [4/12] 类关系图...")
    create_class_diagram(prs)
    print("  [5/12] 功能模块...")
    create_function_modules(prs)
    print("  [6/12] 测试体系...")
    create_test_system(prs)
    print("  [7/12] 单元测试...")
    create_unit_tests(prs)
    print("  [8/12] 集成测试...")
    create_integration_tests(prs)
    print("  [9/12] 覆盖率要求...")
    create_coverage_requirements(prs)
    print("  [10/12] 覆盖率分析...")
    create_coverage_analysis(prs)
    print("  [11/12] 性能测试...")
    create_performance_test(prs)
    print("  [12/12] 总结展望...")
    create_summary(prs)
    
    output_path = "/Users/netchen/Desktop/fossify-math-unit-test/FossifyCalculator_TestReport.pptx"
    prs.save(output_path)
    print(f"\nPPT 已成功生成！")
    print(f"文件路径: {output_path}")
    print(f"总页数: 12 页")

if __name__ == "__main__":
    main()
